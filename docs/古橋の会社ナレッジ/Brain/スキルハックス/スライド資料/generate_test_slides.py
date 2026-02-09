#!/usr/bin/env python3
"""スキルハックス セミナースライド テスト生成（3枚）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design tokens ──
DARK_NAVY = RGBColor(0x1A, 0x1A, 0x2E)
DARK_NAVY_2 = RGBColor(0x12, 0x12, 0x22)
GOLD = RGBColor(0xD4, 0xA8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=DARK_NAVY):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape=MSO_SHAPE.RECTANGLE):
    """Add a shape with optional fill and line."""
    s = slide.shapes.add_shape(shape, left, top, width, height)
    s.line.fill.background()
    if fill_color:
        s.fill.solid()
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    return s


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Arial"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_logo(slide):
    """Add 'Skill Hacks' branding at bottom right."""
    add_text_box(slide,
                 Inches(10.5), Inches(6.8),
                 Inches(2.5), Inches(0.5),
                 "Skill Hacks", font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.RIGHT)


# ============================================================
# Slide 1-01: タイトル画面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide1)

# Spotlight effect: subtle gradient circle in center
spotlight = add_shape(slide1,
                      Inches(3.5), Inches(1.5),
                      Inches(6.3), Inches(4.5),
                      shape=MSO_SHAPE.OVAL)
spotlight.fill.solid()
spotlight.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
spotlight.fill.fore_color.brightness = 0.1

# Top accent line
add_shape(slide1, Inches(4.5), Inches(1.8), Inches(4.3), Pt(2), fill_color=GOLD)

# Main title
add_text_box(slide1,
             Inches(1.5), Inches(2.2),
             Inches(10.3), Inches(2.5),
             'なぜ「勉強熱心な人」ほど\n副業で1円も稼げないのか',
             font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Bottom accent line
add_shape(slide1, Inches(4.5), Inches(4.7), Inches(4.3), Pt(2), fill_color=GOLD)

# Subtitle
add_text_box(slide1,
             Inches(2), Inches(5.2),
             Inches(9.3), Inches(0.5),
             "Skill Hacks 特別動画講座　第1話",
             font_size=16, color=GOLD, bold=False,
             alignment=PP_ALIGN.CENTER)

# Speaker
add_text_box(slide1,
             Inches(2), Inches(5.8),
             Inches(9.3), Inches(0.5),
             "話者：迫 佑樹（さこ ゆうき）",
             font_size=14, color=LIGHT_GRAY, bold=False,
             alignment=PP_ALIGN.CENTER)

add_logo(slide1)


# ============================================================
# Slide 1-02: 3つの価値
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

# Section title
add_text_box(slide2,
             Inches(1), Inches(0.8),
             Inches(11.3), Inches(0.8),
             "この動画シリーズで手に入る3つのこと",
             font_size=32, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Accent line under title
add_shape(slide2, Inches(5), Inches(1.6), Inches(3.3), Pt(2), fill_color=GOLD)

# 3 cards
card_data = [
    ("①", "スキルを最短最速で\n身につける方法", "必要なのは「幹」だけ。\n枝葉は後からいくらでも伸びる。"),
    ("②", "売上を伸ばすための\n知識", "スキルだけでは稼げない。\nお金の流れを理解する。"),
    ("③", "自立した事業者になる\nロードマップ", "副業→本業→資産構築。\n人生を攻略する全体像。"),
]

for i, (num, title, desc) in enumerate(card_data):
    x = Inches(1.2 + i * 3.8)
    y = Inches(2.3)
    w = Inches(3.4)
    h = Inches(4.2)

    # Card background
    card = add_shape(slide2, x, y, w, h, fill_color=RGBColor(0x22, 0x22, 0x38))
    card.line.color.rgb = RGBColor(0x33, 0x33, 0x50)
    card.line.width = Pt(1)

    # Gold number circle
    circle = add_shape(slide2,
                       x + Inches(1.2), y + Inches(0.3),
                       Inches(1.0), Inches(1.0),
                       fill_color=GOLD,
                       shape=MSO_SHAPE.OVAL)
    # Number text on circle
    add_text_box(slide2,
                 x + Inches(1.2), y + Inches(0.45),
                 Inches(1.0), Inches(0.7),
                 num, font_size=28, color=DARK_NAVY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Card title
    add_text_box(slide2,
                 x + Inches(0.2), y + Inches(1.5),
                 Inches(3.0), Inches(1.2),
                 title, font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Card description
    add_text_box(slide2,
                 x + Inches(0.2), y + Inches(2.8),
                 Inches(3.0), Inches(1.2),
                 desc, font_size=12, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

add_logo(slide2)


# ============================================================
# Slide 1-03: 半年後のあなた
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, RGBColor(0x14, 0x14, 0x26))

# Overlay shape for visual depth
overlay = add_shape(slide3,
                    Inches(0), Inches(0),
                    SLIDE_WIDTH, SLIDE_HEIGHT,
                    fill_color=RGBColor(0x1A, 0x1A, 0x2E))
overlay.fill.fore_color.brightness = 0.05

# Top label
add_text_box(slide3,
             Inches(1), Inches(1.2),
             Inches(11.3), Inches(0.6),
             "── FUTURE VISION ──",
             font_size=14, color=GOLD, bold=False,
             alignment=PP_ALIGN.CENTER)

# Main title
add_text_box(slide3,
             Inches(1), Inches(1.8),
             Inches(11.3), Inches(1.0),
             "半年後のあなた",
             font_size=40, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Center accent line
add_shape(slide3, Inches(5.5), Inches(2.9), Inches(2.3), Pt(2), fill_color=GOLD)

# Main catch copy - big gold text
add_text_box(slide3,
             Inches(1.5), Inches(3.5),
             Inches(10.3), Inches(1.5),
             "自分の力で稼いだ、\nという確信",
             font_size=36, color=GOLD, bold=True,
             alignment=PP_ALIGN.CENTER)

# Description bullets
descriptions = [
    "☕  好きなカフェで、ノートPC1台で仕事をしている",
    "💰  月5万円、10万円と副業収入が積み上がっている",
    "🚀  「あのとき始めてよかった」と心から思えている",
]

for i, desc in enumerate(descriptions):
    add_text_box(slide3,
                 Inches(3), Inches(5.0 + i * 0.5),
                 Inches(7.3), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.LEFT)

# Bottom decorative line
add_shape(slide3, Inches(1), Inches(6.8), Inches(11.3), Pt(1), fill_color=RGBColor(0x33, 0x33, 0x50))

add_logo(slide3)


# ============================================================
# Save
# ============================================================
output_path = "/Users/hiroshi/cursor/docs/古橋の会社ナレッジ/Brain/スキルハックス/スライド資料/テスト_第一話スライド.pptx"
prs.save(output_path)
print(f"✅ テストスライドを生成しました（3枚）: {output_path}")
